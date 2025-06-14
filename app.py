import streamlit as st
from PIL import Image
import pptx
import io
import zipfile
import xml.etree.ElementTree as ET
from xml.dom import minidom
import os
import fitz  # Importando a biblioteca PyMuPDF

def criar_manifest_scorm(titulo_curso, nome_recurso_html):
    """
    Cria a estrutura do arquivo imsmanifest.xml para SCORM 1.2.
    """
    manifest = ET.Element("manifest", {
        "identifier": "MANIFEST-1",
        "version": "1.2",
        "xmlns": "http://www.imsproject.org/xsd/imscp_rootv1p1p2",
        "xmlns:adlcp": "http://www.adlnet.org/xsd/adlcp_rootv1p2",
        "xmlns:xsi": "http://www.w3.org/2001/XMLSchema-instance",
        "xsi:schemaLocation": "http://www.imsproject.org/xsd/imscp_rootv1p1p2 imscp_rootv1p1p2.xsd http://www.imsglobal.org/xsd/imsmd_rootv1p2p1 imsmd_rootv1p2p1.xsd http://www.adlnet.org/xsd/adlcp_rootv1p2 adlcp_rootv1p2.xsd"
    })

    metadata = ET.SubElement(manifest, "metadata")
    schema = ET.SubElement(metadata, "schema")
    schema.text = "ADL SCORM"
    schemaversion = ET.SubElement(metadata, "schemaversion")
    schemaversion.text = "1.2"

    organizations = ET.SubElement(manifest, "organizations", {"default": "ORG-1"})
    organization = ET.SubElement(organizations, "organization", {"identifier": "ORG-1", "structure": "hierarchical"})
    title_org = ET.SubElement(organization, "title")
    title_org.text = titulo_curso
    item = ET.SubElement(organization, "item", {"identifier": "ITEM-1", "identifierref": "RES-1", "isvisible": "true"})
    title_item = ET.SubElement(item, "title")
    title_item.text = titulo_curso

    resources = ET.SubElement(manifest, "resources")
    resource = ET.SubElement(resources, "resource", {
        "identifier": "RES-1",
        "type": "webcontent",
        "adlcp:scormtype": "sco",
        "href": nome_recurso_html
    })
    
    file_html = ET.SubElement(resource, "file", {"href": nome_recurso_html})

    xml_str = ET.tostring(manifest, 'utf-8')
    parsed_str = minidom.parseString(xml_str)
    return parsed_str.toprettyxml(indent="  ")


def criar_html_para_conteudo(titulo, corpo_html):
    """
    Cria uma página HTML simples para exibir o conteúdo.
    """
    html_template = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <meta charset="UTF-8">
        <title>{titulo}</title>
        <style>
            body {{ font-family: Arial, sans-serif; margin: 20px; background-color: #f0f2f5; }}
            .container {{ max-width: 900px; margin: auto; background-color: #fff; padding: 20px; box-shadow: 0 0 10px rgba(0,0,0,0.1); border-radius: 8px; }}
            h1 {{ color: #333; }}
            img {{ max-width: 100%; height: auto; display: block; margin: 15px auto; border: 1px solid #ddd; }}
            video {{ max-width: 100%; display: block; margin: 15px auto; }}
        </style>
    </head>
    <body>
        <div class="container">
            <h1>{titulo}</h1>
            <hr>
            {corpo_html}
        </div>
    </body>
    </html>
    """
    return html_template

# --- ALTERAÇÕES DO TÍTULO AQUI ---
st.set_page_config(page_title="Conversor para SCORM - Rovema", layout="wide", initial_sidebar_state="collapsed")
st.title("Conversor para SCORM - Rovema")
# --- FIM DAS ALTERAÇÕES DO TÍTULO ---

st.info("Faça o upload de um arquivo PDF, PPTX ou vídeo (MP4, MOV) para convertê-lo em um pacote SCORM 1.2.")

uploaded_file = st.file_uploader("Selecione um arquivo para conversão", type=["pdf", "pptx", "mp4", "mov"])

if uploaded_file is not None:
    st.markdown("---")
    file_details = {"Nome do Arquivo": uploaded_file.name, "Tipo": uploaded_file.type, "Tamanho (MB)": f"{uploaded_file.size / 1024 / 1024:.2f}"}
    st.write(file_details)

    default_title = os.path.splitext(uploaded_file.name)[0]
    titulo_curso = st.text_input("Título do Curso:", default_title)

    if st.button("✨ Gerar Pacote SCORM"):
        with st.spinner("Processando... Seu pacote SCORM está sendo criado, por favor aguarde."):
            in_memory_zip = io.BytesIO()
            zip_file = zipfile.ZipFile(in_memory_zip, 'w', zipfile.ZIP_DEFLATED)
            corpo_html = ""
            lista_arquivos_recurso = []

            if uploaded_file.type == "application/pdf":
                pdf_doc = fitz.open(stream=uploaded_file.getvalue(), filetype="pdf")
                for i, page in enumerate(pdf_doc):
                    nome_imagem = f"pagina_{i+1}.png"
                    pix = page.get_pixmap(dpi=150) # Aumentar a qualidade da imagem
                    img_bytes = pix.tobytes("png")
                    zip_file.writestr(nome_imagem, img_bytes)
                    lista_arquivos_recurso.append(nome_imagem)
                    corpo_html += f'<img src="{nome_imagem}" alt="Página {i+1}">\n'
                pdf_doc.close()

            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                st.warning("A conversão de PPTX para imagens não é suportada. O texto dos slides será extraído.", icon="⚠️")
                prs = pptx.Presentation(io.BytesIO(uploaded_file.getvalue()))
                for i, slide in enumerate(prs.slides):
                    corpo_html += f"<h3>Slide {i+1}</h3>\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            corpo_html += f"<p>{shape.text}</p>\n"

            elif uploaded_file.type in ["video/mp4", "video/quicktime"]:
                nome_video = uploaded_file.name
                zip_file.writestr(nome_video, uploaded_file.getvalue())
                lista_arquivos_recurso.append(nome_video)
                corpo_html += f'<video controls width="100%"><source src="{nome_video}" type="{uploaded_file.type}"></video>'

            nome_recurso_html = "index.html"
            html_content = criar_html_para_conteudo(titulo_curso, corpo_html)
            zip_file.writestr(nome_recurso_html, html_content)
            
            manifest_content_str = criar_manifest_scorm(titulo_curso, nome_recurso_html)
            
            manifest_root = ET.fromstring(manifest_content_str)
            resource_element = manifest_root.find(".//{http://www.imsproject.org/xsd/imscp_rootv1p1p2}resource[@identifier='RES-1']")
            for file_name in lista_arquivos_recurso:
                 ET.SubElement(resource_element, "file", {"href": file_name})

            final_xml_str = ET.tostring(manifest_root, 'utf-8')
            parsed_str = minidom.parseString(final_xml_str)
            final_manifest_content = parsed_str.toprettyxml(indent="  ")

            zip_file.writestr("imsmanifest.xml", final_manifest_content)
            zip_file.close()
            in_memory_zip.seek(0)

            st.success("Pacote SCORM gerado com sucesso!", icon="✅")

            st.download_button(
                label="📥 Baixar Pacote SCORM (.zip)",
                data=in_memory_zip,
                file_name=f"{titulo_curso.replace(' ', '_')}_scorm.zip",
                mime="application/zip"
            )
